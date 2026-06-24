from __future__ import annotations

import argparse
import io
import json
import os
import sys
import tempfile
from dataclasses import dataclass
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Iterable

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


NAVY = RGBColor(12, 42, 67)
TEAL = RGBColor(0, 104, 121)
GREEN = RGBColor(98, 171, 70)
GANTT_GREEN = RGBColor(112, 173, 71)
GANTT_GREEN_LIGHT = RGBColor(146, 208, 80)
LIGHT_GREEN = RGBColor(202, 229, 190)
LIGHT_BLUE = RGBColor(216, 230, 236)
LIGHT_GRAY = RGBColor(238, 241, 243)
MID_GRAY = RGBColor(168, 175, 180)
DARK_GRAY = RGBColor(67, 73, 77)
WHITE = RGBColor(255, 255, 255)
DESCRIPTION_LIMIT = 145
PROJECT_STATUS_LIMIT = 280

PROJECT_STATUS_FONT_SIZE = 8.0
RAID_HEADER_FONT_SIZE = 10.0
RAID_BODY_FONT_SIZE = 8.7
RAID_DESCRIPTION_TITLE_FONT_SIZE = 8.8
RAID_DESCRIPTION_STATUS_FONT_SIZE = 8.2
GANTT_TITLE_FONT_SIZE = 9.0
GANTT_OWNER_HEADER_FONT_SIZE = 9.0
GANTT_DATE_FONT_SIZE = 8.0
GANTT_TASK_FONT_SIZE = 8.0
GANTT_OWNER_FONT_SIZE = 8.0

# Fixed dashboard geometry inherited from the PowerPoint template.
SUMMARY_LEFT = Inches(0.597)
SUMMARY_TOP = Inches(0.768)
SUMMARY_WIDTH = Inches(4.578)
SUMMARY_HEIGHT = Inches(3.403)
RAID_LEFT = Inches(5.393)
RAID_TOP = SUMMARY_TOP
RAID_WIDTH = Inches(7.35)
RAID_HEIGHT = SUMMARY_HEIGHT
GANTT_LEFT = Inches(0.927)
GANTT_ORIGINAL_TOP = Inches(4.321)
GANTT_TOP = GANTT_ORIGINAL_TOP + Inches(0.20)
GANTT_WIDTH = Inches(11.479)
GANTT_HEIGHT = Inches(1.977)
# Mirror the Cogsdale logo's approximately 0.32-inch right-side spacing.
CLIENT_LOGO_BOX_LEFT = Inches(0.32)
CLIENT_LOGO_BOX_WIDTH = Inches(1.20)
CLIENT_LOGO_BOX_HEIGHT = Inches(0.58)
CLIENT_LOGO_BOTTOM_MARGIN = Inches(0.12)
CLIENT_LOGO_CLEAR_LEFT = Inches(0.30)
CLIENT_LOGO_CLEAR_WIDTH = Inches(1.75)

SUMMARY_COLUMN_WIDTHS = (Inches(1.356), SUMMARY_WIDTH - Inches(1.356))
SUMMARY_ROW_HEIGHTS = (
    Inches(0.295),
    Inches(0.279),
    Inches(0.279),
    Inches(0.279),
    Inches(0.279),
    Inches(0.427),
)
SUMMARY_ROW_HEIGHTS += (SUMMARY_HEIGHT - sum(SUMMARY_ROW_HEIGHTS),)
RAID_COLUMN_WIDTHS = (
    Inches(0.59),
    Inches(3.48),
    Inches(0.78),
    Inches(0.80),
    Inches(0.90),
)
RAID_COLUMN_WIDTHS += (RAID_WIDTH - sum(RAID_COLUMN_WIDTHS),)
RAID_HEADER_HEIGHT = Inches(0.34)
RAID_DATA_ROW_HEIGHT = (RAID_HEIGHT - RAID_HEADER_HEIGHT) // 5
RAID_ROW_HEIGHTS = (
    RAID_HEADER_HEIGHT,
    RAID_DATA_ROW_HEIGHT,
    RAID_DATA_ROW_HEIGHT,
    RAID_DATA_ROW_HEIGHT,
    RAID_DATA_ROW_HEIGHT,
    RAID_HEIGHT - RAID_HEADER_HEIGHT - RAID_DATA_ROW_HEIGHT * 4,
)


@dataclass(frozen=True)
class ProjectSummary:
    project_name: str
    project_manager: str
    project_sponsor: str
    report_date: str
    go_live_date: str
    project_status: str


@dataclass(frozen=True)
class RaidItem:
    item_type: str
    title: str
    latest_update: str
    description: str
    priority: str
    status: str
    assigned: str
    due_date: str
    last_updated: str
    source: str
    reference_link: str
    attachments: str


@dataclass(frozen=True)
class Deliverable:
    name: str
    owner: str
    start: date | None
    end: date | None
    percent_complete: int
    status: str
    level: int | None
    source: str


@dataclass(frozen=True)
class DashboardData:
    project_summary: ProjectSummary
    raid_items: list[RaidItem]
    deliverables: list[Deliverable]


class DashboardDataError(ValueError):
    """Raised when a dashboard JSON file is missing or contains invalid data."""


PROJECT_FIELDS = (
    "project_name",
    "project_manager",
    "project_sponsor",
    "report_date",
    "go_live_date",
    "project_status",
)
def require_object(value: Any, location: str) -> dict[str, Any]:
    if not isinstance(value, dict):
        raise DashboardDataError(f"'{location}' must be a JSON object.")
    return value


def require_list(value: Any, location: str) -> list[Any]:
    if not isinstance(value, list):
        raise DashboardDataError(f"'{location}' must be a JSON array.")
    return value


def require_fields(item: dict[str, Any], fields: tuple[str, ...], location: str) -> None:
    missing = [field for field in fields if field not in item]
    if missing:
        field_list = ", ".join(missing)
        raise DashboardDataError(
            f"'{location}' is missing required field(s): {field_list}."
        )


def require_text(value: Any, location: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise DashboardDataError(f"'{location}' must be a non-empty text value.")
    return value.strip()


def optional_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def limit_description(value: str, limit: int = DESCRIPTION_LIMIT) -> str:
    text = value.strip()
    if len(text) <= limit:
        return text
    if limit <= 1:
        return text[:limit]
    shortened = text[:limit].rsplit(" ", 1)[0].rstrip(" ,;:-")
    if not shortened:
        shortened = text[: limit - 1].rstrip()
    return shortened.rstrip(".") + "."


def require_percent(value: Any, location: str) -> int:
    if isinstance(value, bool) or not isinstance(value, int):
        raise DashboardDataError(
            f"'{location}' must be a whole number between 0 and 100."
        )
    if not 0 <= value <= 100:
        raise DashboardDataError(f"'{location}' must be between 0 and 100.")
    return value


def parse_iso_date(value: Any, location: str) -> date:
    text = require_text(value, location)
    try:
        return date.fromisoformat(text)
    except ValueError as exc:
        raise DashboardDataError(
            f"'{location}' must use YYYY-MM-DD format, for example 2025-10-24."
        ) from exc


def parse_optional_iso_date(value: Any, location: str) -> date | None:
    text = optional_text(value)
    if not text:
        return None
    try:
        return date.fromisoformat(text)
    except ValueError as exc:
        raise DashboardDataError(
            f"'{location}' must use YYYY-MM-DD format, for example 2026-06-15."
        ) from exc


def parse_report_date(value: Any, location: str) -> tuple[str, date]:
    text = require_text(value, location)
    try:
        parsed = datetime.strptime(text, "%B %d, %Y").date()
    except ValueError as exc:
        raise DashboardDataError(
            f"'{location}' must use Month D, YYYY format, for example October 10, 2025."
        ) from exc
    return text, parsed


def dashboard_data_from_dict(raw: Any) -> tuple[DashboardData, date]:
    root = require_object(raw, "dashboard")
    require_fields(root, ("project_summary", "raid_items", "deliverables"), "dashboard")

    project = require_object(root["project_summary"], "project_summary")
    require_fields(project, PROJECT_FIELDS, "project_summary")
    report_date_text, report_date = parse_report_date(
        project["report_date"], "project_summary.report_date"
    )
    summary = ProjectSummary(
        project_name=optional_text(project["project_name"]),
        project_manager=optional_text(project["project_manager"]),
        project_sponsor=optional_text(project["project_sponsor"]),
        report_date=report_date_text,
        go_live_date=optional_text(project["go_live_date"]),
        project_status=optional_text(project["project_status"]),
    )

    raid_items = []
    for index, raw_item in enumerate(require_list(root["raid_items"], "raid_items")):
        location = f"raid_items[{index}]"
        item = require_object(raw_item, location)
        description = limit_description(optional_text(item.get("description")))
        if not description:
            raise DashboardDataError(
                f"'{location}' is missing required field: description."
            )
        item_type = optional_text(item.get("type") or item.get("category"))
        if not item_type:
            raise DashboardDataError(
                f"'{location}' is missing required field: type."
            )
        due_date = optional_text(item.get("due_date") or item.get("due"))
        raid_items.append(
            RaidItem(
                item_type=item_type,
                title=optional_text(item.get("title")),
                latest_update=optional_text(item.get("latest_update")),
                description=description,
                priority=optional_text(item.get("priority")) or "Unassigned",
                status=optional_text(item.get("status")),
                assigned=optional_text(item.get("assigned")),
                due_date=due_date,
                last_updated=optional_text(item.get("last_updated")),
                source=optional_text(item.get("source")),
                reference_link=optional_text(item.get("reference_link")),
                attachments=optional_text(item.get("attachments")),
            )
        )

    deliverables = []
    raw_deliverables = require_list(root["deliverables"], "deliverables")
    if not raw_deliverables:
        raise DashboardDataError(
            "'deliverables' must contain at least one deliverable."
        )
    for index, raw_item in enumerate(raw_deliverables):
        location = f"deliverables[{index}]"
        item = require_object(raw_item, location)
        name = require_text(item.get("name"), f"{location}.name")
        start = parse_optional_iso_date(
            item.get("start_date", item.get("start")),
            f"{location}.start_date",
        )
        end = parse_optional_iso_date(
            item.get("end_date", item.get("end")),
            f"{location}.end_date",
        )
        if start is not None and end is not None and end < start:
            raise DashboardDataError(
                f"'{location}.end_date' cannot be earlier than "
                f"'{location}.start_date'."
            )
        raw_level = item.get("level")
        level = None
        if raw_level not in (None, ""):
            try:
                level = int(raw_level)
            except (TypeError, ValueError) as exc:
                raise DashboardDataError(
                    f"'{location}.level' must be a whole number."
                ) from exc
        deliverables.append(
            Deliverable(
                name=name,
                owner=optional_text(item.get("owner")),
                start=start,
                end=end,
                percent_complete=require_percent(
                    item.get("percent_complete", 0),
                    f"{location}.percent_complete",
                ),
                status=optional_text(item.get("state") or item.get("status"))
                or "Active",
                level=level,
                source=optional_text(item.get("source")),
            )
        )

    return DashboardData(summary, raid_items, deliverables), report_date


def load_dashboard_data(path: Path) -> tuple[DashboardData, date]:
    if not path.exists():
        raise DashboardDataError(f"Data file was not found: {path}")
    if not path.is_file():
        raise DashboardDataError(f"Data path is not a file: {path}")

    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise DashboardDataError(
            f"Could not read '{path}' because the JSON is invalid near line "
            f"{exc.lineno}, column {exc.colno}: {exc.msg}."
        ) from exc
    except OSError as exc:
        raise DashboardDataError(f"Could not read data file '{path}': {exc}") from exc
    return dashboard_data_from_dict(raw)


def _capture_run_style(run) -> dict[str, Any]:
    style = {
        "name": run.font.name,
        "size": run.font.size,
        "bold": run.font.bold,
        "italic": run.font.italic,
        "underline": run.font.underline,
    }
    try:
        style["color_rgb"] = run.font.color.rgb
    except (AttributeError, TypeError):
        style["color_rgb"] = None
    try:
        style["theme_color"] = run.font.color.theme_color
    except (AttributeError, TypeError):
        style["theme_color"] = None
    try:
        style["brightness"] = run.font.color.brightness
    except (AttributeError, TypeError):
        style["brightness"] = None
    return style


def _apply_run_style(run, style: dict[str, Any] | None) -> None:
    if not style:
        return
    run.font.name = style.get("name")
    run.font.size = style.get("size")
    run.font.bold = style.get("bold")
    run.font.italic = style.get("italic")
    run.font.underline = style.get("underline")
    color_applied = False
    if style.get("color_rgb") is not None:
        run.font.color.rgb = style["color_rgb"]
        color_applied = True
    elif (
        style.get("theme_color") is not None
        and getattr(style["theme_color"], "name", "") != "NOT_THEME_COLOR"
    ):
        run.font.color.theme_color = style["theme_color"]
        color_applied = True
    if color_applied and style.get("brightness") is not None:
        run.font.color.brightness = style["brightness"]


def _capture_paragraph_style(paragraph) -> dict[str, Any]:
    return {
        "alignment": paragraph.alignment,
        "level": paragraph.level,
        "line_spacing": paragraph.line_spacing,
        "space_before": paragraph.space_before,
        "space_after": paragraph.space_after,
    }


def _apply_paragraph_style(paragraph, style: dict[str, Any] | None) -> None:
    if not style:
        return
    paragraph.alignment = style["alignment"]
    paragraph.level = style["level"]
    paragraph.line_spacing = style["line_spacing"]
    paragraph.space_before = style["space_before"]
    paragraph.space_after = style["space_after"]


def set_text_preserving_style(text_frame, value: str) -> None:
    source_run_style = None
    source_paragraph_style = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        source_run = text_frame.paragraphs[0].runs[0]
        source_run_style = _capture_run_style(source_run)
        source_paragraph_style = _capture_paragraph_style(text_frame.paragraphs[0])
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    _apply_paragraph_style(paragraph, source_paragraph_style)
    run = paragraph.add_run()
    run.text = value
    _apply_run_style(run, source_run_style)


def set_description_text(text_frame, value: str) -> None:
    value = limit_description(value)
    source_run_style = None
    source_paragraph_style = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        source_run = text_frame.paragraphs[0].runs[0]
        source_run_style = _capture_run_style(source_run)
        source_paragraph_style = _capture_paragraph_style(text_frame.paragraphs[0])

    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    _apply_paragraph_style(paragraph, source_paragraph_style)

    lines = value.split("\n", 1)
    if len(lines) == 2:
        title, status_text = lines
        title_run = paragraph.add_run()
        title_run.text = title.strip()
        _apply_run_style(title_run, source_run_style)
        title_run.font.bold = True
        title_run.font.size = Pt(RAID_DESCRIPTION_TITLE_FONT_SIZE)

        status_run = paragraph.add_run()
        status_run.text = "\n" + status_text.strip()
        _apply_run_style(status_run, source_run_style)
        status_run.font.bold = False
        status_run.font.size = Pt(RAID_DESCRIPTION_STATUS_FONT_SIZE)
    else:
        run = paragraph.add_run()
        run.text = value
        _apply_run_style(run, source_run_style)
        run.font.bold = False
        run.font.size = Pt(RAID_DESCRIPTION_STATUS_FONT_SIZE)


def find_table_shapes(slide):
    summary_shape = None
    raid_shape = None
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        rows = len(shape.table.rows)
        columns = len(shape.table.columns)
        if (rows, columns) == (7, 2):
            summary_shape = shape
        elif columns == 6:
            raid_shape = shape
    if summary_shape is None or raid_shape is None:
        raise ValueError(
            "Template must contain a 7x2 project summary table and a 6-column RAID table."
        )
    return summary_shape, raid_shape


def find_tables(slide):
    summary_shape, raid_shape = find_table_shapes(slide)
    return summary_shape.table, raid_shape.table


def _set_table_geometry(shape, left, top, column_widths, row_heights) -> None:
    shape.left = left
    shape.top = top
    for column, width in zip(shape.table.columns, column_widths):
        column.width = width
    for row, height in zip(shape.table.rows, row_heights):
        row.height = height
    shape.width = sum(column_widths)
    shape.height = sum(row_heights)


def lock_dashboard_layout(summary_shape, raid_shape) -> None:
    _set_table_geometry(
        summary_shape,
        SUMMARY_LEFT,
        SUMMARY_TOP,
        SUMMARY_COLUMN_WIDTHS,
        SUMMARY_ROW_HEIGHTS,
    )
    _set_table_geometry(
        raid_shape,
        RAID_LEFT,
        RAID_TOP,
        RAID_COLUMN_WIDTHS,
        RAID_ROW_HEIGHTS,
    )

    for shape in (summary_shape, raid_shape):
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text_frame.word_wrap = True
                cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE


def _set_text_frame_size(text_frame, size: float) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def replace_placeholders(slide, replacements: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue
        for old, new in replacements.items():
            if old in shape.text:
                set_text_preserving_style(shape.text_frame, shape.text.replace(old, new))


def populate_summary(table, summary: ProjectSummary) -> None:
    project_status = truncate_with_ellipsis(
        summary.project_status, PROJECT_STATUS_LIMIT
    )
    # Combine the former percentage row with Project Status so the table keeps
    # its template footprint while exposing only the six current fields.
    if not table.cell(5, 0).is_merge_origin and not table.cell(5, 0).is_spanned:
        table.cell(5, 0).merge(table.cell(6, 0))
    if not table.cell(5, 1).is_merge_origin and not table.cell(5, 1).is_spanned:
        table.cell(5, 1).merge(table.cell(6, 1))

    labels_and_values = [
        ("Project Name", summary.project_name),
        ("Project Manager", summary.project_manager),
        ("Project Sponsor", summary.project_sponsor),
        ("Date", summary.report_date),
        ("Go-Live Date", summary.go_live_date),
        ("Project Status", project_status),
    ]
    for row_index, (label, value) in enumerate(labels_and_values):
        set_text_preserving_style(table.cell(row_index, 0).text_frame, label)
        set_text_preserving_style(table.cell(row_index, 1).text_frame, value)
        if row_index == 5:
            _set_text_frame_size(
                table.cell(row_index, 1).text_frame,
                PROJECT_STATUS_FONT_SIZE,
            )


def populate_raid(table, items: Iterable[RaidItem]) -> None:
    items = list(items)
    headers = ["RAID", "Description", "Priority", "Status", "Assigned", "Due"]
    for column, header in enumerate(headers):
        set_text_preserving_style(table.cell(0, column).text_frame, header)
        _set_text_frame_size(
            table.cell(0, column).text_frame, RAID_HEADER_FONT_SIZE
        )
        for paragraph in table.cell(0, column).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = WHITE

    available_rows = len(table.rows) - 1
    if len(items) > available_rows:
        raise ValueError(
            f"Template RAID table has room for {available_rows} items; received {len(items)}."
        )

    for row_index in range(1, len(table.rows)):
        values = (
            [
                items[row_index - 1].item_type,
                limit_description(items[row_index - 1].description),
                items[row_index - 1].priority,
                items[row_index - 1].status,
                items[row_index - 1].assigned,
                items[row_index - 1].due_date,
            ]
            if row_index <= len(items)
            else ["", "", "", "", "", ""]
        )
        for column, value in enumerate(values):
            if column == 1 and value:
                set_description_text(table.cell(row_index, column).text_frame, value)
            else:
                set_text_preserving_style(table.cell(row_index, column).text_frame, value)
                _set_text_frame_size(
                    table.cell(row_index, column).text_frame,
                    RAID_BODY_FONT_SIZE if row_index else RAID_HEADER_FONT_SIZE,
                )


def _slide_dimensions(slide) -> tuple[int, int]:
    presentation = slide.part.package.presentation_part.presentation
    return presentation.slide_width, presentation.slide_height


def _validate_geometry(slide, left, top, width, height, label: str) -> None:
    values = (left, top, width, height)
    if any(value is None for value in values):
        raise ValueError(f"{label} has incomplete geometry.")
    if width <= 0 or height <= 0:
        raise ValueError(f"{label} must have positive width and height.")
    slide_width, slide_height = _slide_dimensions(slide)
    tolerance = Inches(0.25)
    if (
        left < -tolerance
        or top < -tolerance
        or left + width > slide_width + tolerance
        or top + height > slide_height + tolerance
    ):
        raise ValueError(f"{label} falls outside the slide bounds.")


def normalize_template_line_geometry(slide) -> None:
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.LINE:
            continue
        if shape.width <= 0:
            shape.width = 1
        if shape.height <= 0:
            shape.height = 1


def _remove_shape_with_relationships(slide, shape) -> None:
    relationship_attributes = {qn("r:id"), qn("r:embed"), qn("r:link")}
    relationship_ids = {
        value
        for element in shape._element.iter()
        for attribute, value in element.attrib.items()
        if attribute in relationship_attributes and value.startswith("rId")
    }
    parent = shape._element.getparent()
    if parent is None:
        return
    parent.remove(shape._element)
    remaining_slide_xml = slide._element.xml
    for relationship_id in relationship_ids:
        if (
            relationship_id not in remaining_slide_xml
            and relationship_id in slide.part.rels
        ):
            slide.part.drop_rel(relationship_id)


def remove_embedded_gantt_safely(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            continue
        _remove_shape_with_relationships(slide, shape)


def clear_embedded_gantt_area(slide) -> None:
    clear_height = GANTT_TOP + GANTT_HEIGHT - GANTT_ORIGINAL_TOP
    add_rect(
        slide,
        GANTT_LEFT,
        GANTT_ORIGINAL_TOP,
        GANTT_WIDTH,
        clear_height,
        WHITE,
    )


def clear_bottom_left_client_logo_area(slide) -> None:
    _, slide_height = _slide_dimensions(slide)
    top = GANTT_TOP + GANTT_HEIGHT
    add_rect(
        slide,
        CLIENT_LOGO_CLEAR_LEFT,
        top,
        CLIENT_LOGO_CLEAR_WIDTH,
        slide_height - top,
        WHITE,
    )


def add_client_logo(
    slide,
    slide_height,
    *,
    client_logo_path: Path | None = None,
    client_logo_bytes: bytes | None = None,
):
    if client_logo_bytes is None and client_logo_path is not None:
        client_logo_bytes = Path(client_logo_path).read_bytes()
    if not client_logo_bytes:
        return None

    try:
        with PILImage.open(io.BytesIO(client_logo_bytes)) as image:
            image_format = (image.format or "").upper()
            pixel_width, pixel_height = image.size
    except (OSError, ValueError) as exc:
        raise ValueError(
            "The client logo could not be read. Use a valid PNG, JPG, or JPEG image."
        ) from exc

    if image_format not in {"PNG", "JPEG"} or pixel_width <= 0 or pixel_height <= 0:
        raise ValueError("The client logo must be a valid PNG, JPG, or JPEG image.")

    scale = min(
        CLIENT_LOGO_BOX_WIDTH / pixel_width,
        CLIENT_LOGO_BOX_HEIGHT / pixel_height,
    )
    logo_width = int(pixel_width * scale)
    logo_height = int(pixel_height * scale)
    box_top = slide_height - CLIENT_LOGO_BOTTOM_MARGIN - CLIENT_LOGO_BOX_HEIGHT
    logo_left = CLIENT_LOGO_BOX_LEFT + (CLIENT_LOGO_BOX_WIDTH - logo_width) // 2
    logo_top = box_top + (CLIENT_LOGO_BOX_HEIGHT - logo_height) // 2
    _validate_geometry(
        slide,
        logo_left,
        logo_top,
        logo_width,
        logo_height,
        "Client logo",
    )

    return slide.shapes.add_picture(
        io.BytesIO(client_logo_bytes),
        logo_left,
        logo_top,
        width=logo_width,
        height=logo_height,
    )


def _is_red_color(color) -> bool:
    if color is None:
        return False
    try:
        red, green, blue = color
    except (TypeError, ValueError):
        return False
    return red >= 150 and green <= 110 and blue <= 110


def remove_gantt_date_markers(slide) -> None:
    gantt_bottom = GANTT_TOP + GANTT_HEIGHT
    for shape in list(slide.shapes):
        if shape.shape_type not in {
            MSO_SHAPE_TYPE.LINE,
            MSO_SHAPE_TYPE.AUTO_SHAPE,
        }:
            continue
        if getattr(shape, "is_placeholder", False):
            continue
        if shape.top < GANTT_TOP - Inches(0.1) or shape.top > gantt_bottom:
            continue
        if (
            shape.width <= 0
            or shape.height < Inches(0.3)
            or shape.width > Pt(5)
            or shape.height <= shape.width * 10
        ):
            continue
        fill_color = None
        line_color = None
        try:
            fill_color = shape.fill.fore_color.rgb
        except (AttributeError, TypeError):
            pass
        try:
            line_color = shape.line.color.rgb
        except (AttributeError, TypeError):
            pass
        if _is_red_color(fill_color) or _is_red_color(line_color):
            _remove_shape_with_relationships(slide, shape)


def add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
    margin: float = 0.04,
):
    _validate_geometry(slide, left, top, width, height, f"Text box '{text[:30]}'")
    shape = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_rect(slide, left, top, width, height, fill, line=None):
    _validate_geometry(slide, left, top, width, height, "Rectangle")
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    return shape


def calculate_gantt_timeline(
    deliverables: list[Deliverable],
    report_date: date,
    week_count: int = 9,
) -> tuple[date, date, list[date], str]:
    del deliverables
    timeline_start = report_date - timedelta(days=report_date.weekday())
    labels = [
        timeline_start + timedelta(days=7 * index)
        for index in range(week_count)
    ]
    timeline_end = timeline_start + timedelta(days=7 * week_count)
    return timeline_start, timeline_end, labels, "weekly"


def truncate_with_ellipsis(value: str, max_characters: int) -> str:
    text = " ".join(value.split())
    if len(text) <= max_characters:
        return text
    if max_characters <= 3:
        return "." * max_characters
    shortened = text[: max_characters - 3].rsplit(" ", 1)[0].rstrip()
    if not shortened:
        shortened = text[: max_characters - 3].rstrip()
    return shortened + "..."


def clip_deliverable_range(
    start: date | None,
    end: date | None,
    timeline_start: date,
    timeline_end: date,
) -> tuple[date, date] | None:
    if start is None or end is None:
        return None
    if start >= timeline_end or end < timeline_start:
        return None
    visible_start = max(start, timeline_start)
    visible_end = min(end, timeline_end)
    if visible_end < visible_start:
        return None
    return visible_start, visible_end


def draw_gantt(slide, deliverables: list[Deliverable], report_date: date) -> None:
    left = GANTT_LEFT
    top = GANTT_TOP
    width = GANTT_WIDTH
    height = GANTT_HEIGHT
    name_column_width = Inches(3.40)
    owner_width = Inches(1.35)
    timeline_left = left + name_column_width + owner_width
    timeline_width = width - name_column_width - owner_width
    header_height = Inches(0.36)
    row_height = (height - header_height) / len(deliverables)

    chart_start, chart_end, timeline_labels, _ = calculate_gantt_timeline(
        deliverables, report_date
    )
    total_days = max(1, (chart_end - chart_start).days)

    add_rect(slide, left, top, width, height, WHITE, MID_GRAY)
    add_text(
        slide,
        "Current & Upcoming Deliverables",
        left,
        top,
        name_column_width,
        header_height,
        size=GANTT_TITLE_FONT_SIZE,
        color=WHITE,
        bold=True,
        fill=NAVY,
    )
    add_text(
        slide,
        "Owner",
        left + name_column_width,
        top,
        owner_width,
        header_height,
        size=GANTT_OWNER_HEADER_FONT_SIZE,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=NAVY,
    )

    weekly_column_width = timeline_width / len(timeline_labels)
    for label_index, label_date in enumerate(timeline_labels):
        x = timeline_left + weekly_column_width * label_index
        add_text(
            slide,
            label_date.strftime("%b %d"),
            x,
            top,
            weekly_column_width,
            header_height,
            size=GANTT_DATE_FONT_SIZE,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            fill=TEAL if label_index % 2 == 0 else NAVY,
            margin=0,
        )

    for index, deliverable in enumerate(deliverables):
        y = top + header_height + index * row_height
        row_fill = WHITE if index % 2 == 0 else LIGHT_GRAY
        add_rect(slide, left, y, width, row_height, row_fill)
        add_text(
            slide,
            truncate_with_ellipsis(deliverable.name, 40),
            left,
            y,
            name_column_width,
            row_height,
            size=GANTT_TASK_FONT_SIZE,
            color=DARK_GRAY,
            bold=index == 0,
        )
        add_text(
            slide,
            truncate_with_ellipsis(deliverable.owner, 16),
            left + name_column_width,
            y,
            owner_width,
            row_height,
            size=GANTT_OWNER_FONT_SIZE,
            color=DARK_GRAY,
            align=PP_ALIGN.CENTER,
        )

        visible_range = clip_deliverable_range(
            deliverable.start,
            deliverable.end,
            chart_start,
            chart_end,
        )
        if visible_range is None:
            continue
        visible_start, visible_end = visible_range
        start_fraction = (visible_start - chart_start).days / total_days
        duration_fraction = max(1, (visible_end - visible_start).days) / total_days
        bar_left = timeline_left + timeline_width * start_fraction
        bar_width = timeline_width * duration_fraction
        bar_top = y + row_height * 0.25
        bar_height = row_height * 0.5
        if bar_width <= 0 or bar_height <= 0:
            continue
        base_color = GANTT_GREEN_LIGHT
        progress_color = GANTT_GREEN
        add_rect(slide, bar_left, bar_top, bar_width, bar_height, base_color)
        if deliverable.percent_complete:
            add_rect(
                slide,
                bar_left,
                bar_top,
                max(Inches(0.04), bar_width * deliverable.percent_complete / 100),
                bar_height,
                progress_color,
            )

def create_dashboard_presentation(
    template: Path,
    dashboard_data: DashboardData,
    report_date: date,
    *,
    client_logo_path: Path | None = None,
    client_logo_bytes: bytes | None = None,
):
    if not template.exists():
        raise ValueError(f"Template file was not found: {template}")
    presentation = Presentation(template)
    if not presentation.slides:
        raise ValueError("Template contains no slides.")

    slide = presentation.slides[0]
    normalize_template_line_geometry(slide)
    replace_placeholders(
        slide,
        {
            "{{TITLE}}": "Weekly Project Status Report",
            "Weekly Project Status Report": "Weekly Project Status Report",
        },
    )
    summary_shape, raid_shape = find_table_shapes(slide)
    lock_dashboard_layout(summary_shape, raid_shape)
    populate_summary(summary_shape.table, dashboard_data.project_summary)
    populate_raid(raid_shape.table, dashboard_data.raid_items)
    clear_bottom_left_client_logo_area(slide)
    add_client_logo(
        slide,
        presentation.slide_height,
        client_logo_path=client_logo_path,
        client_logo_bytes=client_logo_bytes,
    )
    remove_embedded_gantt_safely(slide)
    clear_embedded_gantt_area(slide)
    remove_gantt_date_markers(slide)
    draw_gantt(slide, dashboard_data.deliverables, report_date)
    return presentation


def validate_generated_pptx(path: str | Path) -> list[str]:
    warnings: list[str] = []
    pptx_path = Path(path)
    if not pptx_path.exists():
        return [f"Generated PowerPoint file does not exist: {pptx_path}"]
    if not pptx_path.is_file():
        return [f"Generated PowerPoint path is not a file: {pptx_path}"]
    if pptx_path.stat().st_size <= 0:
        return [f"Generated PowerPoint file is empty: {pptx_path}"]

    try:
        presentation = Presentation(pptx_path)
    except (OSError, ValueError, KeyError) as exc:
        return [f"Generated PowerPoint could not be reopened: {exc}"]

    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    tolerance = Inches(0.25)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            shape_name = getattr(shape, "name", f"shape {shape.shape_id}")
            if shape.width <= 0 or shape.height <= 0:
                warnings.append(
                    f"Slide {slide_index} '{shape_name}' has zero or negative dimensions."
                )
                continue
            if (
                shape.left < -tolerance
                or shape.top < -tolerance
                or shape.left + shape.width > slide_width + tolerance
                or shape.top + shape.height > slide_height + tolerance
            ):
                warnings.append(
                    f"Slide {slide_index} '{shape_name}' extends outside the slide bounds."
                )
    return warnings


def _save_presentation_safely(presentation, output: Path) -> None:
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    handle = tempfile.NamedTemporaryFile(
        mode="wb",
        suffix=".pptx",
        prefix=f".{output.stem}-",
        dir=output.parent,
        delete=False,
    )
    temporary_path = Path(handle.name)
    handle.close()
    try:
        presentation.save(temporary_path)
        warnings = validate_generated_pptx(temporary_path)
        if warnings:
            raise ValueError(
                "Generated PowerPoint failed validation: " + " ".join(warnings)
            )
        os.replace(temporary_path, output)
    finally:
        temporary_path.unlink(missing_ok=True)


def build_dashboard_bytes(
    template: Path,
    dashboard_data: DashboardData,
    report_date: date,
    *,
    client_logo_path: Path | None = None,
    client_logo_bytes: bytes | None = None,
) -> bytes:
    presentation = create_dashboard_presentation(
        template,
        dashboard_data,
        report_date,
        client_logo_path=client_logo_path,
        client_logo_bytes=client_logo_bytes,
    )
    with tempfile.TemporaryDirectory(prefix="weekly-dashboard-") as directory:
        output = Path(directory) / "weekly-project-status-report.pptx"
        _save_presentation_safely(presentation, output)
        return output.read_bytes()


def build_dashboard(
    template: Path,
    output: Path,
    dashboard_data: DashboardData,
    report_date: date,
    *,
    client_logo_path: Path | None = None,
    client_logo_bytes: bytes | None = None,
) -> None:
    presentation = create_dashboard_presentation(
        template,
        dashboard_data,
        report_date,
        client_logo_path=client_logo_path,
        client_logo_bytes=client_logo_bytes,
    )

    _save_presentation_safely(presentation, output)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a one-slide weekly project dashboard from a PowerPoint template."
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=Path("templates/weekly-status-template.pptx"),
        help="Path to the source PowerPoint template.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("generated/weekly-project-status-report.pptx"),
        help="Path for the generated dashboard.",
    )
    parser.add_argument(
        "--data",
        type=Path,
        default=Path("sample_inputs/dashboard_sample.json"),
        help=(
            "Path to the dashboard JSON file. Defaults to "
            "sample_inputs/dashboard_sample.json."
        ),
    )
    parser.add_argument(
        "--client-logo",
        type=Path,
        default=None,
        help="Optional PNG, JPG, or JPEG client logo for the bottom-left footer.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    arguments = parse_args()
    try:
        data, data_report_date = load_dashboard_data(arguments.data)
        build_dashboard(
            arguments.template,
            arguments.output,
            data,
            data_report_date,
            client_logo_path=arguments.client_logo,
        )
    except (DashboardDataError, ValueError, OSError) as error:
        print(f"Error: {error}", file=sys.stderr)
        sys.exit(2)
    print(f"Dashboard created successfully: {arguments.output}")
