import copy
import io
import json
import tempfile
import unittest
import zipfile
from datetime import date
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from generate_dashboard import (
    CLIENT_LOGO_BOTTOM_MARGIN,
    CLIENT_LOGO_BOX_HEIGHT,
    CLIENT_LOGO_BOX_LEFT,
    CLIENT_LOGO_BOX_WIDTH,
    GANTT_DATE_FONT_SIZE,
    GANTT_GREEN,
    GANTT_GREEN_LIGHT,
    GANTT_HEIGHT,
    GANTT_ORIGINAL_TOP,
    GANTT_OWNER_HEADER_FONT_SIZE,
    GANTT_OWNER_FONT_SIZE,
    GANTT_TASK_FONT_SIZE,
    GANTT_TITLE_FONT_SIZE,
    GANTT_TOP,
    PROJECT_STATUS_LIMIT,
    PROJECT_STATUS_FONT_SIZE,
    RAID_BODY_FONT_SIZE,
    RAID_COLUMN_WIDTHS,
    RAID_DESCRIPTION_STATUS_FONT_SIZE,
    RAID_DESCRIPTION_TITLE_FONT_SIZE,
    RAID_HEADER_FONT_SIZE,
    RAID_ROW_HEIGHTS,
    SUMMARY_COLUMN_WIDTHS,
    SUMMARY_ROW_HEIGHTS,
    build_dashboard,
    build_dashboard_bytes,
    create_dashboard_presentation,
    dashboard_data_from_dict,
    find_table_shapes,
    truncate_with_ellipsis,
    validate_generated_pptx,
)


TEMPLATE_PATH = Path("templates/weekly-status-template.pptx")
SAMPLE_PATH = Path("sample_inputs/dashboard_sample.json")


def make_logo_bytes(image_format: str, size: tuple[int, int]) -> bytes:
    output = io.BytesIO()
    mode = "RGBA" if image_format == "PNG" else "RGB"
    color = (24, 113, 146, 180) if mode == "RGBA" else (24, 113, 146)
    Image.new(mode, size, color).save(output, format=image_format)
    return output.getvalue()


def bottom_left_pictures(presentation):
    return [
        shape
        for shape in presentation.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape.left < CLIENT_LOGO_BOX_LEFT + CLIENT_LOGO_BOX_WIDTH
        and shape.top >= GANTT_TOP + GANTT_HEIGHT
    ]


def dashboard_data(project_status="On track.", raid_count=1, long_raid=False):
    raw = json.loads(SAMPLE_PATH.read_text(encoding="utf-8"))
    raw["project_summary"]["project_status"] = project_status
    source_items = raw["raid_items"] or [
        {
            "type": "Action",
            "description": "Test item\nStatus: Work is progressing.",
            "priority": "Normal",
            "status": "In Progress",
            "assigned": "PM",
            "due_date": "2026-07-01",
        }
    ]
    items = []
    for index in range(raid_count):
        item = copy.deepcopy(source_items[index % len(source_items)])
        if long_raid:
            item["description"] = (
                f"Long RAID Item {index + 1}\nStatus: "
                + "This is a deliberately long project update used to verify fixed "
                "table geometry and description truncation. " * 8
            )
        items.append(item)
    raw["raid_items"] = items
    return dashboard_data_from_dict(raw)


def gantt_title_top(presentation):
    slide = presentation.slides[0]
    shape = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text == "Current & Upcoming Deliverables"
    )
    return shape.top


class PowerPointLayoutTests(unittest.TestCase):
    def test_raid_top_aligns_with_summary_and_geometry_is_fixed(self):
        data, report_date = dashboard_data(raid_count=5, long_raid=True)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        summary_shape, raid_shape = find_table_shapes(presentation.slides[0])

        self.assertEqual(summary_shape.top, raid_shape.top)
        self.assertEqual(
            [column.width for column in summary_shape.table.columns],
            list(SUMMARY_COLUMN_WIDTHS),
        )
        self.assertEqual(
            [row.height for row in summary_shape.table.rows],
            list(SUMMARY_ROW_HEIGHTS),
        )
        self.assertEqual(
            [column.width for column in raid_shape.table.columns],
            list(RAID_COLUMN_WIDTHS),
        )
        self.assertEqual(
            [row.height for row in raid_shape.table.rows],
            list(RAID_ROW_HEIGHTS),
        )

    def test_title_style_remains_visible_after_safe_text_replacement(self):
        data, report_date = dashboard_data(raid_count=1)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        title = next(
            shape
            for shape in presentation.slides[0].shapes
            if shape.has_text_frame and shape.text == "Weekly Project Status Report"
        )
        run = title.text_frame.paragraphs[0].runs[0]

        self.assertEqual(run.font.color.brightness, -0.5)

    def test_gantt_top_is_fixed_for_status_length_and_raid_count(self):
        short_data, short_date = dashboard_data("On track.", raid_count=1)
        long_data, long_date = dashboard_data(
            "A long project status narrative describing progress, risks, decisions, "
            "testing, validation, configuration, and next steps. " * 8,
            raid_count=5,
            long_raid=True,
        )

        short_presentation = create_dashboard_presentation(
            TEMPLATE_PATH, short_data, short_date
        )
        long_presentation = create_dashboard_presentation(
            TEMPLATE_PATH, long_data, long_date
        )

        self.assertEqual(gantt_title_top(short_presentation), GANTT_TOP)
        self.assertEqual(gantt_title_top(long_presentation), GANTT_TOP)

    def test_long_project_status_is_truncated_without_moving_tables(self):
        data, report_date = dashboard_data("Project status detail " * 80, raid_count=3)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        summary_shape, raid_shape = find_table_shapes(presentation.slides[0])
        status_text = summary_shape.table.cell(5, 1).text

        self.assertLessEqual(len(status_text), PROJECT_STATUS_LIMIT)
        self.assertTrue(status_text.endswith("..."))
        status_run = summary_shape.table.cell(5, 1).text_frame.paragraphs[0].runs[0]
        self.assertEqual(status_run.font.size.pt, PROJECT_STATUS_FONT_SIZE)
        self.assertEqual(summary_shape.top, raid_shape.top)
        self.assertEqual(gantt_title_top(presentation), GANTT_TOP)

    def test_generated_output_removes_template_red_gantt_marker(self):
        template = Presentation(TEMPLATE_PATH)
        slide = template.slides[0]
        marker = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(8),
            GANTT_TOP + Inches(0.36),
            Pt(1.4),
            Inches(1.4),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(201, 48, 44)
        marker.line.fill.background()

        with tempfile.TemporaryDirectory() as directory:
            modified_template = Path(directory) / "template-with-marker.pptx"
            template.save(modified_template)
            data, report_date = dashboard_data(raid_count=5, long_raid=True)
            output = build_dashboard_bytes(modified_template, data, report_date)

        reopened = Presentation(io.BytesIO(output))
        self.assertEqual(len(reopened.slides), 1)
        for shape in reopened.slides[0].shapes:
            if shape.width > Pt(5) or shape.height < Inches(0.3):
                continue
            try:
                self.assertNotEqual(shape.fill.fore_color.rgb, RGBColor(201, 48, 44))
            except (AttributeError, TypeError):
                continue

    def test_raid_font_sizes_are_larger_and_description_style_is_preserved(self):
        data, report_date = dashboard_data(raid_count=5, long_raid=True)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        _, raid_shape = find_table_shapes(presentation.slides[0])
        table = raid_shape.table

        header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
        description_runs = table.cell(1, 1).text_frame.paragraphs[0].runs

        self.assertEqual(header_run.font.size.pt, RAID_HEADER_FONT_SIZE)
        self.assertGreaterEqual(RAID_HEADER_FONT_SIZE, 10.0)
        self.assertTrue(header_run.font.bold)
        self.assertEqual(header_run.font.color.rgb, RGBColor(255, 255, 255))
        self.assertAlmostEqual(body_run.font.size.pt, RAID_BODY_FONT_SIZE, places=1)
        self.assertGreater(RAID_BODY_FONT_SIZE, 8.0)
        self.assertAlmostEqual(
            description_runs[0].font.size.pt,
            RAID_DESCRIPTION_TITLE_FONT_SIZE,
            places=1,
        )
        self.assertTrue(description_runs[0].font.bold)
        self.assertGreater(RAID_DESCRIPTION_TITLE_FONT_SIZE, 8.0)
        if len(description_runs) > 1:
            self.assertAlmostEqual(
                description_runs[1].font.size.pt,
                RAID_DESCRIPTION_STATUS_FONT_SIZE,
                places=1,
            )
            self.assertFalse(description_runs[1].font.bold)
            self.assertGreater(RAID_DESCRIPTION_STATUS_FONT_SIZE, 7.6)

    def test_gantt_fonts_green_bars_and_lower_position(self):
        data, report_date = dashboard_data(raid_count=5)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        slide = presentation.slides[0]

        task_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == data.deliverables[0].name
        )
        owner_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame
            and shape.text == truncate_with_ellipsis(data.deliverables[0].owner, 16)
        )
        date_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == "Oct 06"
        )
        title_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame
            and shape.text == "Current & Upcoming Deliverables"
        )
        owner_header_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == "Owner"
        )
        colors = []
        for shape in slide.shapes:
            try:
                colors.append(shape.fill.fore_color.rgb)
            except (AttributeError, TypeError):
                continue

        self.assertEqual(
            task_shape.text_frame.paragraphs[0].runs[0].font.size.pt,
            GANTT_TASK_FONT_SIZE,
        )
        self.assertEqual(
            owner_shape.text_frame.paragraphs[0].runs[0].font.size.pt,
            GANTT_OWNER_FONT_SIZE,
        )
        self.assertEqual(
            date_shape.text_frame.paragraphs[0].runs[0].font.size.pt,
            GANTT_DATE_FONT_SIZE,
        )
        self.assertEqual(
            title_shape.text_frame.paragraphs[0].runs[0].font.size.pt,
            GANTT_TITLE_FONT_SIZE,
        )
        self.assertEqual(
            owner_header_shape.text_frame.paragraphs[0].runs[0].font.size.pt,
            GANTT_OWNER_HEADER_FONT_SIZE,
        )
        self.assertGreater(GANTT_TASK_FONT_SIZE, 7.3)
        self.assertGreater(GANTT_OWNER_FONT_SIZE, 7.2)
        self.assertGreater(GANTT_DATE_FONT_SIZE, 7.0)
        self.assertIn(GANTT_GREEN_LIGHT, colors)
        self.assertIn(GANTT_GREEN, colors)
        self.assertEqual(GANTT_TOP - GANTT_ORIGINAL_TOP, Inches(0.20))
        self.assertLessEqual(GANTT_TOP + GANTT_HEIGHT, presentation.slide_height)

    def test_blank_project_information_stays_blank_in_powerpoint(self):
        raw = json.loads(SAMPLE_PATH.read_text(encoding="utf-8"))
        raw["project_summary"].update(
            {
                "project_name": "",
                "project_manager": "",
                "project_sponsor": "",
                "go_live_date": "",
                "project_status": "",
            }
        )
        data, report_date = dashboard_data_from_dict(raw)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        summary_shape, _ = find_table_shapes(presentation.slides[0])

        for row_index in (0, 1, 2, 4, 5):
            self.assertEqual(summary_shape.table.cell(row_index, 1).text, "")

    def test_percent_complete_is_removed_and_summary_uses_six_visible_fields(self):
        data, report_date = dashboard_data(raid_count=1)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        summary_shape, _ = find_table_shapes(presentation.slides[0])
        table = summary_shape.table

        self.assertEqual(
            [table.cell(row, 0).text for row in range(6)],
            [
                "Project Name",
                "Project Manager",
                "Project Sponsor",
                "Date",
                "Go-Live Date",
                "Project Status",
            ],
        )
        self.assertTrue(table.cell(5, 0).is_merge_origin)
        self.assertTrue(table.cell(6, 0).is_spanned)
        all_text = "\n".join(
            shape.text
            for shape in presentation.slides[0].shapes
            if shape.has_text_frame
        )
        self.assertNotIn("% Complete, Per WIP Hours", all_text)

    def test_no_client_logo_leaves_bottom_left_blank_and_right_logo_remains(self):
        data, report_date = dashboard_data(raid_count=5)
        presentation = create_dashboard_presentation(
            TEMPLATE_PATH, data, report_date
        )
        pictures = [
            shape
            for shape in presentation.slides[0].shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]

        self.assertEqual(bottom_left_pictures(presentation), [])
        self.assertTrue(any(shape.left > Inches(10.0) for shape in pictures))

        output = io.BytesIO()
        presentation.save(output)
        reopened = Presentation(io.BytesIO(output.getvalue()))
        self.assertEqual(len(reopened.slides), 1)

    def test_png_client_logo_is_scaled_proportionally_into_bottom_left(self):
        data, report_date = dashboard_data(raid_count=5)
        logo_bytes = make_logo_bytes("PNG", (240, 100))
        output = build_dashboard_bytes(
            TEMPLATE_PATH,
            data,
            report_date,
            client_logo_bytes=logo_bytes,
        )
        presentation = Presentation(io.BytesIO(output))
        logo = bottom_left_pictures(presentation)[-1]
        box_top = (
            presentation.slide_height
            - CLIENT_LOGO_BOTTOM_MARGIN
            - CLIENT_LOGO_BOX_HEIGHT
        )

        self.assertAlmostEqual(logo.width / logo.height, 2.4, places=2)
        self.assertLess(CLIENT_LOGO_BOX_LEFT, Inches(0.55))
        self.assertAlmostEqual(
            CLIENT_LOGO_BOX_LEFT / Inches(1),
            0.32,
            places=2,
        )
        self.assertGreaterEqual(logo.left, CLIENT_LOGO_BOX_LEFT)
        self.assertLessEqual(
            logo.left + logo.width,
            CLIENT_LOGO_BOX_LEFT + CLIENT_LOGO_BOX_WIDTH,
        )
        self.assertGreaterEqual(logo.top, box_top)
        self.assertLessEqual(
            logo.top + logo.height,
            presentation.slide_height - CLIENT_LOGO_BOTTOM_MARGIN,
        )
        self.assertGreaterEqual(logo.top, GANTT_TOP + GANTT_HEIGHT)
        self.assertGreaterEqual(logo.left, 0)
        self.assertGreaterEqual(logo.top, 0)
        self.assertLessEqual(
            logo.left + logo.width,
            presentation.slide_width,
        )
        self.assertLessEqual(
            logo.top + logo.height,
            presentation.slide_height,
        )
        self.assertTrue(
            any(
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and shape.left > Inches(10.0)
                for shape in presentation.slides[0].shapes
            )
        )

    def test_jpeg_client_logo_path_is_supported_and_aspect_ratio_is_preserved(self):
        data, report_date = dashboard_data(raid_count=5)
        with tempfile.TemporaryDirectory() as directory:
            logo_path = Path(directory) / "client-logo.jpg"
            logo_path.write_bytes(make_logo_bytes("JPEG", (100, 220)))
            presentation = create_dashboard_presentation(
                TEMPLATE_PATH,
                data,
                report_date,
                client_logo_path=logo_path,
            )

        logo = bottom_left_pictures(presentation)[-1]
        self.assertAlmostEqual(logo.width / logo.height, 100 / 220, places=2)
        self.assertLessEqual(logo.width, CLIENT_LOGO_BOX_WIDTH)
        self.assertLessEqual(logo.height, CLIENT_LOGO_BOX_HEIGHT)

    def test_saved_powerpoint_reopens_validates_and_has_positive_shape_dimensions(self):
        data, report_date = dashboard_data(raid_count=5, long_raid=True)
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "validated-dashboard.pptx"
            build_dashboard(TEMPLATE_PATH, output, data, report_date)

            self.assertTrue(output.exists())
            self.assertGreater(output.stat().st_size, 0)
            self.assertEqual(validate_generated_pptx(output), [])
            reopened = Presentation(output)

        self.assertEqual(len(reopened.slides), 1)
        self.assertTrue(
            all(
                shape.width > 0 and shape.height > 0
                for shape in reopened.slides[0].shapes
            )
        )

    def test_embedded_template_gantt_and_relationships_are_removed_cleanly(self):
        data, report_date = dashboard_data(raid_count=5)
        output = build_dashboard_bytes(TEMPLATE_PATH, data, report_date)

        presentation = Presentation(io.BytesIO(output))
        self.assertFalse(
            any(
                shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
                for shape in presentation.slides[0].shapes
            )
        )
        with zipfile.ZipFile(io.BytesIO(output)) as package:
            names = package.namelist()
            slide_relationships = package.read(
                "ppt/slides/_rels/slide1.xml.rels"
            ).decode("utf-8")

        self.assertFalse(any(name.startswith("ppt/embeddings/") for name in names))
        self.assertNotIn("relationships/package", slide_relationships)

    def test_validator_reports_missing_and_empty_files(self):
        with tempfile.TemporaryDirectory() as directory:
            missing = Path(directory) / "missing.pptx"
            empty = Path(directory) / "empty.pptx"
            empty.write_bytes(b"")

            self.assertTrue(validate_generated_pptx(missing))
            self.assertTrue(validate_generated_pptx(empty))


if __name__ == "__main__":
    unittest.main()
