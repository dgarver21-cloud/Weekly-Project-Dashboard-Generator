import unittest
from datetime import date, timedelta
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from generate_dashboard import (
    Deliverable,
    GANTT_HEIGHT,
    GANTT_TOP,
    calculate_gantt_timeline,
    clip_deliverable_range,
    dashboard_data_from_dict,
    draw_gantt,
    populate_summary,
    truncate_with_ellipsis,
)


def deliverable(name, start, end, owner="Owner"):
    return Deliverable(
        name=name,
        owner=owner,
        start=start,
        end=end,
        percent_complete=0,
        status="Active",
        level=2,
        source="Test",
    )


class GanttTimelineTests(unittest.TestCase):
    def test_report_date_starts_on_monday_of_report_week(self):
        start, end, labels, mode = calculate_gantt_timeline(
            [], date(2025, 10, 10)
        )
        self.assertEqual(start, date(2025, 10, 6))
        self.assertEqual(end, date(2025, 12, 8))
        self.assertEqual(mode, "weekly")

    def test_exactly_nine_weekly_headers_are_created(self):
        _, _, labels, _ = calculate_gantt_timeline([], date(2025, 10, 10))
        self.assertEqual(len(labels), 9)
        self.assertEqual(
            labels,
            [
                date(2025, 10, 6),
                date(2025, 10, 13),
                date(2025, 10, 20),
                date(2025, 10, 27),
                date(2025, 11, 3),
                date(2025, 11, 10),
                date(2025, 11, 17),
                date(2025, 11, 24),
                date(2025, 12, 1),
            ],
        )
        for first, second in zip(labels, labels[1:]):
            self.assertEqual(second - first, timedelta(days=7))

    def test_bar_is_clipped_at_left_edge(self):
        clipped = clip_deliverable_range(
            date(2025, 9, 1),
            date(2025, 10, 20),
            date(2025, 10, 6),
            date(2025, 12, 8),
        )
        self.assertEqual(clipped, (date(2025, 10, 6), date(2025, 10, 20)))

    def test_bar_is_clipped_at_right_edge(self):
        clipped = clip_deliverable_range(
            date(2025, 11, 20),
            date(2026, 1, 31),
            date(2025, 10, 6),
            date(2025, 12, 8),
        )
        self.assertEqual(clipped, (date(2025, 11, 20), date(2025, 12, 8)))

    def test_non_overlapping_and_missing_dates_have_no_bar(self):
        self.assertIsNone(
            clip_deliverable_range(
                date(2026, 1, 1),
                date(2026, 1, 10),
                date(2025, 10, 6),
                date(2025, 12, 8),
            )
        )
        self.assertIsNone(
            clip_deliverable_range(
                None,
                date(2025, 11, 1),
                date(2025, 10, 6),
                date(2025, 12, 8),
            )
        )

    def test_long_name_and_owner_are_truncated_cleanly(self):
        long_name = (
            "Business Central Configuration and Integration Validation Activities"
        )
        long_owner = "Danvers and Cogsdale Implementation Leadership Team"
        shortened_name = truncate_with_ellipsis(long_name, 43)
        shortened_owner = truncate_with_ellipsis(long_owner, 18)
        self.assertLessEqual(len(shortened_name), 43)
        self.assertLessEqual(len(shortened_owner), 18)
        self.assertTrue(shortened_name.endswith("..."))
        self.assertTrue(shortened_owner.endswith("..."))
        self.assertFalse(shortened_name.endswith(" ..."))

    def test_powerpoint_summary_uses_selected_report_date(self):
        dashboard_json = json.loads(
            Path("sample_inputs/dashboard_sample.json").read_text(encoding="utf-8")
        )
        dashboard_json["project_summary"]["report_date"] = "June 18, 2026"
        dashboard_data, parsed_date = dashboard_data_from_dict(dashboard_json)

        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(7, 2, 0, 0, 4000000, 3000000).table
        populate_summary(table, dashboard_data.project_summary)

        self.assertEqual(parsed_date, date(2026, 6, 18))
        self.assertEqual(table.cell(3, 1).text, "June 18, 2026")

    def test_gantt_draws_without_red_report_date_line(self):
        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        items = [
            deliverable(
                "Business Central Configuration",
                date(2026, 6, 15),
                date(2026, 7, 10),
            )
        ]

        draw_gantt(slide, items, date(2026, 6, 18))

        red_marker = RGBColor(201, 48, 44)
        shape_colors = []
        for shape in slide.shapes:
            try:
                shape_colors.append(shape.fill.fore_color.rgb)
            except (AttributeError, TypeError):
                continue
        self.assertNotIn(red_marker, shape_colors)
        self.assertIn("Jun 15", [shape.text for shape in slide.shapes if shape.has_text_frame])
        self.assertTrue(all(shape.top >= GANTT_TOP for shape in slide.shapes))
        self.assertTrue(
            all(shape.top + shape.height <= GANTT_TOP + GANTT_HEIGHT for shape in slide.shapes)
        )


if __name__ == "__main__":
    unittest.main()
