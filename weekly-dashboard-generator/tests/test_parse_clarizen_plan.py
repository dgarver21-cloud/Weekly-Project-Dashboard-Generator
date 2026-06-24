import unittest
from datetime import date, datetime
from io import BytesIO

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from generate_dashboard import Deliverable, draw_gantt

from parse_clarizen_plan import (
    ClarizenPlanError,
    parse_clarizen_workbook,
    rank_deliverables,
)


def workbook_bytes(rows, sheet_name="Work Plan"):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    for row in rows:
        sheet.append(row)
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


class ParseClarizenPlanTests(unittest.TestCase):
    def test_resource_is_used_before_owner(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Due Date", "Level", "Owner", "Resource"],
                ["Active", "Configuration", datetime(2026, 6, 1), datetime(2026, 6, 20), 2, "Dylan Garver", "Alex Smith"],
            ]
        )

        deliverables, _ = parse_clarizen_workbook(data)
        self.assertEqual(deliverables[0]["owner"], "Alex Smith")

        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        item = deliverables[0]
        draw_gantt(
            slide,
            [
                Deliverable(
                    name=item["name"],
                    owner=item["owner"],
                    start=date.fromisoformat(item["start_date"]),
                    end=date.fromisoformat(item["end_date"]),
                    percent_complete=0,
                    status=item["state"],
                    level=item["level"],
                    source=item["source"],
                )
            ],
            date(2026, 6, 10),
        )
        gantt_text = [shape.text for shape in slide.shapes if shape.has_text_frame]
        self.assertIn("Alex Smith", gantt_text)
        self.assertNotIn("Dylan Garver", gantt_text)

    def test_owner_is_used_when_resource_is_blank(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Due Date", "Level", "Owner", "Resource"],
                ["Active", "Configuration", datetime(2026, 6, 1), datetime(2026, 6, 20), 2, "Dylan Garver", ""],
            ]
        )

        deliverables, _ = parse_clarizen_workbook(data)
        self.assertEqual(deliverables[0]["owner"], "Dylan Garver")

    def test_blank_resource_and_owner_remain_blank(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Due Date", "Level", "Owner", "Resource"],
                ["Active", "Configuration", datetime(2026, 6, 1), datetime(2026, 6, 20), 2, "", ""],
            ]
        )

        deliverables, warnings = parse_clarizen_workbook(data)
        self.assertEqual(deliverables[0]["owner"], "")
        self.assertTrue(any("owner/resource" in warning for warning in warnings))

    def test_dynamic_header_filtering_due_date_and_warnings(self):
        data = workbook_bytes(
            [
                ["Clarizen Project Plan Export"],
                ["Generated for weekly reporting"],
                [],
                [
                    " State ",
                    " Name ",
                    " Start Date ",
                    " Due Date ",
                    " Level ",
                    " Owner ",
                ],
                ["Active", "Portfolio", datetime(2026, 6, 1), datetime(2026, 6, 5), 0, "PM"],
                ["Active", "Phase", datetime(2026, 6, 1), datetime(2026, 6, 5), 1, "PM"],
                ["Active", "Included Task", datetime(2026, 6, 2), datetime(2026, 6, 20), 2, "Alex"],
                ["Completed", "Completed Task", datetime(2026, 6, 2), datetime(2026, 6, 20), 2, "Alex"],
                ["Active", "Missing Dates", None, None, 3, ""],
            ]
        )

        deliverables, warnings = parse_clarizen_workbook(
            data, today=date(2026, 6, 10)
        )

        self.assertEqual(
            [item["name"] for item in deliverables],
            ["Included Task", "Missing Dates"],
        )
        self.assertEqual(deliverables[0]["start_date"], "2026-06-02")
        self.assertEqual(deliverables[0]["end_date"], "2026-06-20")
        self.assertEqual(deliverables[0]["level"], 2)
        self.assertEqual(
            deliverables[0]["ranking_reason"],
            "Currently active",
        )
        self.assertTrue(any("missing start date" in warning for warning in warnings))
        self.assertTrue(any("missing due/end date" in warning for warning in warnings))
        self.assertTrue(any("missing an owner" in warning for warning in warnings))

    def test_excel_serial_dates_are_converted(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "End Date", "Level"],
                ["Active", "Serial Date Task", 46174, 46180, 2],
            ]
        )
        deliverables, _ = parse_clarizen_workbook(data, today=date(2026, 1, 1))
        self.assertRegex(deliverables[0]["start_date"], r"^\d{4}-\d{2}-\d{2}$")
        self.assertRegex(deliverables[0]["end_date"], r"^\d{4}-\d{2}-\d{2}$")

    def test_missing_end_date_column_is_rejected(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Level"],
                ["Active", "Task", datetime(2026, 6, 1), 2],
            ]
        )
        with self.assertRaisesRegex(ClarizenPlanError, "Due Date or End Date"):
            parse_clarizen_workbook(data)

    def test_invoicing_milestones_and_subtree_are_excluded(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Due Date", "Level", "Owner"],
                [
                    "Active",
                    "Normal Active Deliverable",
                    datetime(2026, 6, 1),
                    datetime(2026, 6, 20),
                    2,
                    "Alex",
                ],
                [
                    "Active",
                    "Project Invoicing",
                    datetime(2026, 6, 1),
                    datetime(2026, 6, 30),
                    2,
                    "PM",
                ],
                [
                    "Active",
                    "M-1 Project Invoice",
                    datetime(2026, 6, 5),
                    datetime(2026, 6, 10),
                    3,
                    "PM",
                ],
                [
                    "Active",
                    "MS-1 Invoice Milestone",
                    datetime(2026, 6, 10),
                    datetime(2026, 6, 15),
                    3,
                    "PM",
                ],
                [
                    "Active",
                    "Billing Review",
                    datetime(2026, 6, 15),
                    datetime(2026, 6, 20),
                    3,
                    "PM",
                ],
                [
                    "Active",
                    "Deliverable After Invoicing",
                    datetime(2026, 6, 20),
                    datetime(2026, 7, 1),
                    2,
                    "Alex",
                ],
            ]
        )

        deliverables, warnings = parse_clarizen_workbook(
            data, today=date(2026, 6, 9)
        )
        names = [item["name"] for item in deliverables]

        self.assertEqual(
            names,
            ["Normal Active Deliverable", "Deliverable After Invoicing"],
        )
        self.assertTrue(
            any("Excluded 4 Project Invoicing" in warning for warning in warnings)
        )

    def test_standalone_m_and_ms_milestones_are_excluded_case_insensitively(self):
        data = workbook_bytes(
            [
                ["State", "Name", "Start Date", "Due Date", "Level"],
                ["Active", "m-1 Project Invoice", None, None, 2],
                ["Active", "ms-10 Billing Milestone", None, None, 2],
                ["Active", "Migrate Data", None, None, 2],
            ]
        )

        deliverables, _ = parse_clarizen_workbook(data)
        self.assertEqual([item["name"] for item in deliverables], ["Migrate Data"])

    def test_current_and_upcoming_rank_before_old_and_admin_tasks(self):
        deliverables = [
            {
                "name": "FIS Project Management",
                "start_date": "2026-01-01",
                "end_date": "2026-12-31",
            },
            {
                "name": "Old Active Task",
                "start_date": "2025-01-01",
                "end_date": "2025-12-31",
            },
            {
                "name": "Current Configuration",
                "start_date": "2026-06-01",
                "end_date": "2026-06-20",
            },
            {
                "name": "Upcoming Training",
                "start_date": "2026-06-20",
                "end_date": "2026-07-15",
            },
            {
                "name": "Missing Date Task",
                "start_date": "",
                "end_date": "",
            },
            {
                "name": "Ending Soon",
                "start_date": "2026-01-01",
                "end_date": "2026-07-01",
            },
        ]
        ranked = rank_deliverables(
            deliverables, reference_date=date(2026, 6, 10), recommendation_limit=5
        )
        names = [item["name"] for item in ranked]
        self.assertEqual(names[0], "Current Configuration")
        self.assertIn("Upcoming Training", names[:4])
        self.assertEqual(names[-1], "FIS Project Management")
        admin = next(
            item for item in ranked if item["name"] == "FIS Project Management"
        )
        self.assertEqual(admin["ranking_reason"], "Long-running/admin task")
        self.assertFalse(admin["recommended"])


if __name__ == "__main__":
    unittest.main()
